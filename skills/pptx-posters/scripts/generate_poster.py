#!/usr/bin/env python3
"""Generate one macro-free PPTX poster from approved, strictly local JSON."""

from __future__ import annotations

import argparse
import os
import sys
from datetime import datetime
from importlib.metadata import PackageNotFoundError, version
from pathlib import Path
from typing import Any

from _common import (
    CliError,
    checked_output_file,
    commit_temp_file,
    emit_json,
    private_temp_file,
    sha256_file,
)
from _manifest import load_and_validate_manifest
from _pptx import (
    analyze_layout,
    patch_accessibility_metadata,
    require_safe_pptx,
    strip_generated_printer_settings,
)
from inventory_images import build_inventory

EXPECTED_VERSIONS = {
    "python-pptx": "1.0.2",
    "Pillow": "12.3.0",
    "lxml": "6.1.1",
}


def _require_exact_dependencies() -> dict[str, str]:
    installed: dict[str, str] = {}
    for distribution, expected in EXPECTED_VERSIONS.items():
        try:
            actual = version(distribution)
        except PackageNotFoundError as exc:
            raise CliError(
                f"{distribution}=={expected} is required; install exact pins with "
                '`uv pip install "python-pptx==1.0.2" "Pillow==12.3.0" '
                '"lxml==6.1.1"`'
            ) from exc
        if actual != expected:
            raise CliError(
                f"{distribution}=={expected} is required for reproducible "
                f"generation, but {actual} is installed"
            )
        installed[distribution] = actual
    return installed


def _rgb(value: str, rgb_class: Any) -> Any:
    return rgb_class(
        int(value[1:3], 16),
        int(value[3:5], 16),
        int(value[5:7], 16),
    )


def _set_text_shape(
    slide: Any,
    element: dict[str, Any],
    *,
    order: int,
    colors: dict[str, str],
    contrast_pairs: dict[str, dict[str, Any]],
    Inches: Any,
    Pt: Any,
    RGBColor: Any,
    PP_ALIGN: Any,
    MSO_ANCHOR: Any,
    MSO_AUTO_SIZE: Any,
    existing_shape: Any | None = None,
) -> None:
    if existing_shape is None:
        shape = slide.shapes.add_textbox(
            Inches(float(element["x_in"])),
            Inches(float(element["y_in"])),
            Inches(float(element["width_in"])),
            Inches(float(element["height_in"])),
        )
    else:
        shape = existing_shape
        shape.left = Inches(float(element["x_in"]))
        shape.top = Inches(float(element["y_in"]))
        shape.width = Inches(float(element["width_in"]))
        shape.height = Inches(float(element["height_in"]))
    shape.name = f"R{order:03d}_TEXT_{element['id']}"
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.text = element["text"]
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    margin = Inches(float(element["margin_in"]))
    text_frame.margin_left = margin
    text_frame.margin_right = margin
    text_frame.margin_top = margin
    text_frame.margin_bottom = margin
    text_frame.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }[element["vertical_align"]]

    pair = contrast_pairs[element["contrast_pair_id"]]
    foreground = colors[pair["foreground_color_id"]]
    background = colors[pair["background_color_id"]]
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(background, RGBColor)
    line_color_id = element["line_color_id"]
    if line_color_id is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = _rgb(colors[line_color_id], RGBColor)
        shape.line.width = Pt(float(element["line_width_pt"]))

    alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }[element["align"]]
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = alignment
        for run in paragraph.runs:
            run.font.name = element["font_face"]
            run.font.size = Pt(float(element["font_size_pt_design"]))
            run.font.bold = bool(element["bold"])
            run.font.color.rgb = _rgb(foreground, RGBColor)


def _add_picture(
    slide: Any,
    element: dict[str, Any],
    *,
    order: int,
    asset: dict[str, Any],
    asset_path: Path,
    metadata: dict[str, Any],
    Inches: Any,
) -> tuple[str, str]:
    width_px = int(metadata["width_px"])
    height_px = int(metadata["height_px"])
    box_width = float(element["width_in"])
    box_height = float(element["height_in"])
    fit = min(box_width / width_px, box_height / height_px)
    placed_width = width_px * fit
    placed_height = height_px * fit
    left = float(element["x_in"]) + (box_width - placed_width) / 2.0
    top = float(element["y_in"]) + (box_height - placed_height) / 2.0
    picture = slide.shapes.add_picture(
        str(asset_path),
        Inches(left),
        Inches(top),
        width=Inches(placed_width),
        height=Inches(placed_height),
    )
    shape_name = f"R{order:03d}_IMAGE_{element['id']}"
    picture.name = shape_name
    return shape_name, asset["alt_text"]


def _build_presentation(
    document: dict[str, Any],
    validation: dict[str, Any],
    *,
    image_inventory: dict[str, Any],
    output_path: Path,
) -> dict[str, str]:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise CliError(
            "python-pptx is required; install the exact pins with "
            '`uv pip install "python-pptx==1.0.2" "Pillow==12.3.0" '
            '"lxml==6.1.1"`'
        ) from exc

    presentation = Presentation()
    presentation.slide_width = Inches(float(document["canvas"]["width_in"]))
    presentation.slide_height = Inches(float(document["canvas"]["height_in"]))
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    title_placeholder = slide.shapes.title
    if title_placeholder is None:
        raise CliError("built-in title-only layout did not provide a title placeholder")
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(
        document["canvas"]["background_color"], RGBColor
    )

    core = presentation.core_properties
    core.title = document["document"]["title"]
    core.subject = document["document"]["subject"]
    core.author = "; ".join(document["document"]["authors"])
    core.language = document["document"]["language"]
    approved_at = document["approval"]["approved_at"]
    normalized_time = (
        approved_at[:-1] + "+00:00" if approved_at.endswith("Z") else approved_at
    )
    approval_time = datetime.fromisoformat(normalized_time)
    core.created = approval_time
    core.modified = approval_time

    colors = {
        color_id: value.upper()
        for color_id, value in document["palette"]["colors"].items()
    }
    pairs = {
        pair["id"]: pair for pair in document["palette"]["contrast_pairs"]
    }
    assets = {asset["id"]: asset for asset in document["assets"]}
    paths = {
        asset_id: Path(path)
        for asset_id, path in validation["asset_paths"].items()
    }
    metadata = {
        record["id"]: record["metadata"] for record in image_inventory["images"]
    }
    alt_text_by_shape_name: dict[str, str] = {}
    for element in document["elements"]:
        order = int(element["reading_order"])
        if element["type"] == "text":
            _set_text_shape(
                slide,
                element,
                order=order,
                colors=colors,
                contrast_pairs=pairs,
                Inches=Inches,
                Pt=Pt,
                RGBColor=RGBColor,
                PP_ALIGN=PP_ALIGN,
                MSO_ANCHOR=MSO_ANCHOR,
                MSO_AUTO_SIZE=MSO_AUTO_SIZE,
                existing_shape=(
                    title_placeholder if element["role"] == "title" else None
                ),
            )
        else:
            asset = assets[element["asset_id"]]
            shape_name, alt_text = _add_picture(
                slide,
                element,
                order=order,
                asset=asset,
                asset_path=paths[element["asset_id"]],
                metadata=metadata[element["asset_id"]],
                Inches=Inches,
            )
            alt_text_by_shape_name[shape_name] = alt_text
    presentation.save(output_path)
    return alt_text_by_shape_name


def generate_poster(
    manifest_value: str | Path,
    output_value: str | Path,
) -> dict[str, Any]:
    """Generate, patch, technically inspect, and publish a new PPTX poster."""
    dependency_versions = _require_exact_dependencies()
    destination = checked_output_file(output_value, suffix=".pptx")
    manifest_path, document, validation = load_and_validate_manifest(
        manifest_value,
        verify_assets=True,
        require_approval=True,
    )
    image_inventory = build_inventory(manifest_path, document, validation)
    if not image_inventory["pass"]:
        messages = "; ".join(issue["message"] for issue in image_inventory["issues"])
        raise CliError(f"image inventory failed: {messages}")

    source_descriptor, source_temp = private_temp_file(
        destination, suffix=".pptx"
    )
    sanitized_descriptor, sanitized_temp = private_temp_file(
        destination, suffix=".pptx"
    )
    patched_descriptor, patched_temp = private_temp_file(
        destination, suffix=".pptx"
    )
    os.close(source_descriptor)
    os.close(sanitized_descriptor)
    os.close(patched_descriptor)
    try:
        alt_texts = _build_presentation(
            document,
            validation,
            image_inventory=image_inventory,
            output_path=source_temp,
        )
        strip_generated_printer_settings(source_temp, sanitized_temp)
        require_safe_pptx(sanitized_temp)
        patch_accessibility_metadata(
            sanitized_temp,
            patched_temp,
            alt_text_by_shape_name=alt_texts,
            language=document["document"]["language"],
        )
        package_report = require_safe_pptx(patched_temp)
        accessibility = package_report["accessibility"]
        if accessibility["pictures_missing_alt_text"]:
            raise CliError("generated PPTX contains a picture without alt text")
        if accessibility["slide_title_count"] != 1:
            raise CliError(
                "generated PPTX must contain exactly one native slide-title placeholder"
            )
        if accessibility["text_runs_missing_language"]:
            raise CliError("generated PPTX contains text without explicit language")
        layout_report = analyze_layout(
            patched_temp,
            print_scale=float(validation["physical_output"]["print_scale"]),
            minimum_font_pt_final=float(
                document["quality"]["minimum_font_pt_final"]
            ),
        )
        expected_order = [
            (
                f"R{int(element['reading_order']):03d}_"
                f"{'TEXT' if element['type'] == 'text' else 'IMAGE'}_"
                f"{element['id']}"
            )
            for element in document["elements"]
        ]
        actual_order = [
            item["name"]
            for item in layout_report["slides"][0]["reading_order"]
        ]
        if actual_order != expected_order:
            raise CliError(
                "generated PPTX direct shape order does not match approved reading order"
            )
        if not layout_report["pass"]:
            codes = ", ".join(issue["code"] for issue in layout_report["issues"])
            raise CliError(f"generated PPTX failed layout checks: {codes}")
        commit_temp_file(patched_temp, destination)
    finally:
        for temporary in (source_temp, sanitized_temp, patched_temp):
            try:
                temporary.unlink()
            except FileNotFoundError:
                pass

    return {
        "schema_version": "1.0",
        "generated": True,
        "output": {
            "path": str(destination),
            "sha256": sha256_file(destination),
            "size_bytes": destination.stat().st_size,
        },
        "manifest": {
            "path": str(manifest_path),
            "content_sha256": validation["content_sha256"],
            "approved_by": document["approval"]["approved_by"],
            "approved_at": document["approval"]["approved_at"],
        },
        "dependencies": dependency_versions,
        "package_security": {
            "safe": package_report["safe"],
            "member_count": package_report["package"]["member_count"],
            "external_relationships": 0,
            "macros_or_embedded_objects": 0,
        },
        "technical_checks": {
            "layout_pass": layout_report["pass"],
            "image_inventory_pass": image_inventory["pass"],
            "pictures_with_alt_text": accessibility["pictures_with_alt_text"],
            "picture_count": accessibility["picture_count"],
            "native_slide_title_count": accessibility["slide_title_count"],
            "text_runs_with_language": accessibility["text_runs_with_language"],
            "text_run_count": accessibility["text_run_count"],
        },
        "manual_checks_required": [
            "PowerPoint Accessibility Checker and Reading Order pane",
            "alt-text and native long-description completeness",
            "screen-reader and keyboard navigation",
            "rendered text overflow, font substitution, equations, and glyphs",
            "PDF page size, tags, font rendering, and image quality after export",
            "physical color proof, trim, bleed, safe margin, and final scaling",
            "QR scan test plus visible fallback URL/text",
            "author sign-off on every claim, number, citation, and source",
        ],
    }


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description=(
            "Generate one editable, macro-free PPTX poster from strict, approved "
            "local JSON and local PNG/JPEG assets. The command refuses placeholders, "
            "unknown sources, unapproved content, remote assets, unsafe package "
            "features, low effective DPI, overlap, bounds errors, and overwrites."
        )
    )
    parser.add_argument("manifest", help="approved local poster manifest JSON")
    parser.add_argument(
        "--output",
        required=True,
        help="new .pptx path; an existing destination is never overwritten",
    )
    parser.add_argument("--report", help="optional new JSON generation report path")
    return parser


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    try:
        args = parser.parse_args(argv)
        if args.report is not None:
            checked_output_file(args.report, suffix=".json")
        report = generate_poster(args.manifest, args.output)
        emit_json(report, output=args.report)
        return 0
    except CliError as exc:
        parser.exit(2, f"error: {exc}\n")


if __name__ == "__main__":
    sys.exit(main())
