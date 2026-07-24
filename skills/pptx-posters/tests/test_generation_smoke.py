"""Exact-pinned synthetic PPTX generation and audit smoke test."""

from __future__ import annotations

import hashlib
import json
import stat
import sys
import tempfile
import unittest
import xml.etree.ElementTree as ET
import zipfile
from importlib.metadata import PackageNotFoundError, version
from pathlib import Path

SKILL_ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = SKILL_ROOT / "scripts"
TESTS = SKILL_ROOT / "tests"
sys.path.insert(0, str(SCRIPTS))
sys.path.insert(0, str(TESTS))

from _common import CliError  # noqa: E402
from _manifest import load_and_validate_manifest, manifest_content_hash  # noqa: E402
from _pptx import P_NS, analyze_layout, inspect_pptx  # noqa: E402
from check_layout import apply_manifest_checks  # noqa: E402
from generate_poster import generate_poster  # noqa: E402
from inventory_images import build_inventory  # noqa: E402
from synthetic import build_manifest  # noqa: E402

EXPECTED = {
    "python-pptx": "1.0.2",
    "Pillow": "12.3.0",
    "lxml": "6.1.1",
}


def exact_dependencies_available() -> bool:
    try:
        return all(version(name) == expected for name, expected in EXPECTED.items())
    except PackageNotFoundError:
        return False


@unittest.skipUnless(
    exact_dependencies_available(),
    "run with the exact pins in assets/generation_dependencies.json",
)
class GenerationSmokeTests(unittest.TestCase):
    def test_hidden_image_metadata_is_a_release_blocker(self) -> None:
        from PIL import Image, PngImagePlugin

        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            image_path = root / "synthetic.png"
            metadata = PngImagePlugin.PngInfo()
            metadata.add_text("Comment", "Synthetic hidden metadata")
            Image.new("RGB", (600, 480), "white").save(
                image_path,
                format="PNG",
                pnginfo=metadata,
            )
            digest = hashlib.sha256(image_path.read_bytes()).hexdigest()
            manifest = build_manifest(
                manifest_content_hash,
                include_image=True,
                image_sha256=digest,
            )
            validation = {
                "manifest_path": str(root / "poster.json"),
                "content_sha256": manifest["approval"]["content_sha256"],
                "physical_output": {"print_scale": 1.0},
            }
            inventory = build_inventory(
                root / "poster.json",
                manifest,
                validation,
            )
        self.assertFalse(inventory["pass"])
        self.assertIn(
            "HIDDEN_IMAGE_METADATA",
            {issue["code"] for issue in inventory["issues"]},
        )

    def test_generate_inspect_layout_inventory_alt_text_and_no_overwrite(self) -> None:
        from PIL import Image, ImageDraw
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            image_path = root / "synthetic.png"
            image = Image.new("RGB", (600, 480), "white")
            drawing = ImageDraw.Draw(image)
            drawing.rectangle((60, 60, 540, 420), fill="#4477AA")
            image.save(image_path, format="PNG")
            digest = hashlib.sha256(image_path.read_bytes()).hexdigest()

            manifest = build_manifest(
                manifest_content_hash,
                include_image=True,
                image_sha256=digest,
            )
            manifest_path = root / "poster.json"
            manifest_path.write_text(
                json.dumps(manifest, indent=2) + "\n",
                encoding="utf-8",
            )
            output = root / "poster.pptx"
            generation = generate_poster(manifest_path, output)

            self.assertTrue(output.is_file())
            self.assertEqual(
                stat.S_IMODE(output.stat().st_mode),
                0o600,
            )
            self.assertEqual(generation["dependencies"], EXPECTED)
            self.assertEqual(generation["technical_checks"]["picture_count"], 1)
            self.assertEqual(
                generation["technical_checks"]["pictures_with_alt_text"],
                1,
            )

            package = inspect_pptx(output)
            self.assertTrue(package["safe"], package["findings"])
            self.assertEqual(
                package["accessibility"]["pictures_missing_alt_text"],
                [],
            )
            self.assertEqual(package["accessibility"]["slide_title_count"], 1)
            self.assertEqual(
                package["accessibility"]["text_runs_missing_language"],
                [],
            )

            layout = analyze_layout(
                output,
                print_scale=1.0,
                minimum_font_pt_final=18.0,
            )
            self.assertTrue(layout["pass"], layout["issues"])
            self.assertEqual(layout["slide_size"]["width_in"], 10.0)
            self.assertEqual(layout["slide_size"]["height_in"], 8.0)
            _, approved_document, approved_validation = load_and_validate_manifest(
                manifest_path
            )
            apply_manifest_checks(
                layout,
                approved_document,
                approved_validation,
            )
            self.assertEqual(layout["issues"], [])

            validation = {
                "manifest_path": str(manifest_path),
                "content_sha256": manifest["approval"]["content_sha256"],
                "physical_output": {"print_scale": 1.0},
            }
            inventory = build_inventory(manifest_path, manifest, validation)
            self.assertTrue(inventory["pass"], inventory["issues"])
            placement = inventory["images"][0]["placements"][0]
            self.assertGreaterEqual(
                placement["minimum_effective_dpi_final"],
                100.0,
            )

            with zipfile.ZipFile(output, "r") as archive:
                self.assertFalse(
                    any(
                        name.startswith("ppt/printerSettings/")
                        for name in archive.namelist()
                    )
                )
                slide = ET.fromstring(archive.read("ppt/slides/slide1.xml"))
            descriptions = [
                node.attrib.get("descr")
                for node in slide.findall(f".//{{{P_NS}}}cNvPr")
                if node.attrib.get("name", "").endswith("_IMAGE_image")
            ]
            self.assertEqual(
                descriptions,
                [
                    "Blue rectangle on white used only to verify local image placement."
                ],
            )
            reopened = Presentation(output)
            self.assertEqual(len(reopened.slides), 1)
            self.assertEqual(len(reopened.slides[0].shapes), 3)
            self.assertEqual(
                reopened.slides[0].shapes.title.text,
                "Synthetic layout test",
            )

            second_output = root / "poster-second.pptx"
            second_generation = generate_poster(manifest_path, second_output)
            self.assertEqual(
                generation["output"]["sha256"],
                second_generation["output"]["sha256"],
            )

            with self.assertRaisesRegex(CliError, "overwrite"):
                generate_poster(manifest_path, output)


if __name__ == "__main__":
    unittest.main()
